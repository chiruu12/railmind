"""Generate the Rail Saarthi pitch deck (docs/Rail_Saarthi_Pitch.pptx).

Branded 16:9 deck matching the logo palette (dark navy + electric blue +
amber). Eight content slides (problem, pain points, solution, control-room
demo, passenger app, architecture, impact, roadmap) plus a cover. Run:

    python3 scripts/make_ppt.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


def _strip_theme_style(shape) -> None:
    """Drop the <p:style> element so a shape renders with ONLY its explicit
    fill/line — no inherited theme drop-shadow, fill, line or font color.
    Without this, autoshapes pick up the Office theme's effectRef (a muddy
    drop shadow) which looks terrible on the dark deck."""
    el = shape._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)

ROOT = Path(__file__).resolve().parent.parent
BRAND = ROOT / "frontend" / "public" / "brand"
OUT = ROOT / "docs" / "Rail_Saarthi_Pitch.pptx"

# ── palette ───────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x01, 0x07, 0x15)       # background (matches logo art bg)
PANEL = RGBColor(0x12, 0x1B, 0x30)      # card fill
PANEL2 = RGBColor(0x16, 0x22, 0x3C)     # raised card
LINE = RGBColor(0x24, 0x33, 0x52)       # borders
WHITE = RGBColor(0xF4, 0xF7, 0xFC)
MUTE = RGBColor(0x9A, 0xA7, 0xBD)       # secondary text
BLUE = RGBColor(0x4A, 0x9E, 0xFF)       # accent blue (Saarthi)
AMBER = RGBColor(0xFF, 0x9E, 0x1B)      # accent amber (orbit)
GREEN = RGBColor(0x35, 0xD0, 0x7F)

# 16:9
EMUW = Emu(12192000)
EMUH = Emu(6858000)


def _in(v: float) -> Emu:
    return Emu(int(v * 914400))


def slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    _strip_theme_style(bg)
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def textbox(s, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    return tf


def run(p, text, *, size, color=WHITE, bold=False, italic=False, font="Arial", spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return r


def para(tf, *, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    return p


def rect(s, x, y, w, h, *, fill=PANEL, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    sh = s.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    _strip_theme_style(sh)
    try:
        sh.adjustments[0] = radius
    except (IndexError, KeyError):
        pass
    return sh


def chip(s, x, y, text, color=BLUE):
    """Small uppercase eyebrow label with an accent bar."""
    bar = rect(s, x, y + 0.04, 0.32, 0.07, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, x + 0.45, y - 0.07, 6.0, 0.4)
    p = para(tf, first=True)
    run(p, text.upper(), size=12, color=color, bold=True, spacing=2.5)
    return bar


def accent_underline(s, x, y, w):
    rect(s, x, y, w * 0.6, 0.045, fill=AMBER, line=None, shape=MSO_SHAPE.RECTANGLE)
    rect(s, x + w * 0.6, y, w * 0.4, 0.045, fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)


def page_no(s, n):
    tf = textbox(s, 11.3, 7.0, 0.8, 0.3, align=PP_ALIGN.RIGHT)
    run(para(tf, first=True), f"{n:02d}", size=10, color=MUTE, bold=True, spacing=1)
    tf2 = textbox(s, 0.6, 7.0, 4, 0.3)
    run(para(tf2, first=True), "Rail Saarthi", size=10, color=MUTE, bold=True, spacing=1)


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = EMUW
prs.slide_height = EMUH


# ── Cover ─────────────────────────────────────────────────────────────────────
s = slide(prs)
# subtle top accent
rect(s, 0, 0, 13.333, 0.12, fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
lockup = BRAND / "rail-saarthi-lockup.png"
if lockup.exists():
    pic = s.shapes.add_picture(str(lockup), _in(3.55), _in(1.55), height=_in(2.5))
tf = textbox(s, 1.5, 4.35, 10.333, 1.2, align=PP_ALIGN.CENTER)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "An agentic operating system for Indian Railways", size=22, color=WHITE, bold=True)
p2 = para(tf, align=PP_ALIGN.CENTER, space_before=8)
run(p2, "Multi-agent AI that keeps trains on time, in real time.", size=15, color=MUTE)
accent_underline(s, 5.4, 5.65, 2.5)
tf3 = textbox(s, 1.5, 6.6, 10.333, 0.5, align=PP_ALIGN.CENTER)
run(para(tf3, first=True, align=PP_ALIGN.CENTER),
    "FAR AWAY 2026  ·  Hackathon MVP  ·  built on Hive", size=12, color=MUTE, spacing=1)


# ── Slide 1 — Problem ─────────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "The Problem")
tf = textbox(s, 0.7, 1.05, 11.9, 1.6)
p = para(tf, first=True, line=1.05)
run(p, "India's railways are falling behind, ", size=40, color=WHITE, bold=True)
run(p, "on time.", size=40, color=AMBER, bold=True)

# big stat block — vertical comparison (older above, current below)
rect(s, 0.7, 3.0, 6.0, 3.2, fill=PANEL, line=LINE)
tf = textbox(s, 1.1, 3.4, 5.4, 0.4)
run(para(tf, first=True), "ON-TIME PERFORMANCE, KEY ROUTES", size=13, color=MUTE, bold=True, spacing=1.5)
tf = textbox(s, 1.1, 4.0, 5.4, 0.7)
p = para(tf, first=True)
run(p, "94.2%", size=32, color=MUTE, bold=True)
run(p, "   in 2020", size=14, color=MUTE)
tf = textbox(s, 1.1, 4.7, 5.4, 1.0)
p = para(tf, first=True)
run(p, "73.6%", size=54, color=AMBER, bold=True)
run(p, "   by 2023", size=14, color=MUTE)
tf = textbox(s, 1.1, 5.72, 5.4, 0.4)
run(para(tf, first=True), "A 20.6 point fall in just three years.", size=13, color=MUTE, italic=True)

# supporting facts
facts = [
    ("19M", "passengers carried every single day"),
    ("69,000 km", "of track across the network"),
    ("0", "real-time, system-wide intelligence today"),
]
fy = 3.0
for big, small in facts:
    rect(s, 7.0, fy, 5.6, 0.95, fill=PANEL, line=LINE)
    tf = textbox(s, 7.35, fy + 0.12, 2.0, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True), big, size=26, color=BLUE, bold=True)
    tf = textbox(s, 9.3, fy + 0.12, 3.1, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, line=1.0), small, size=13, color=WHITE)
    fy += 1.13
page_no(s, 1)


# ── Slide 2 — Pain points ─────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Current Pain Points")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "Where the system breaks down today", size=32, color=WHITE, bold=True)

cards = [
    ("Chronic delays", "80% of key routes run over capacity, with no slack to absorb a single slip."),
    ("Fragmented scheduling", "Crew and platform assignments juggled by hand, control center to control center."),
    ("Aging infrastructure", "British-era signaling still gates throughput on critical sections."),
    ("No proactive comms", "Passengers learn of disruptions late, if at all, with no live, trustworthy updates."),
    ("Freight vs passenger conflict", "Shared tracks force constant trade-offs with no system-wide optimizer."),
    ("Siloed control centers", "Each zone sees its own slice; no shared, real-time network awareness."),
]
cw, ch, gx, gy = 3.83, 1.95, 0.2, 0.25
x0, y0 = 0.7, 2.15
for i, (title, body) in enumerate(cards):
    col, row = i % 3, i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=PANEL, line=LINE)
    rect(s, x, y, 0.09, ch, fill=AMBER if row else BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, x + 0.3, y + 0.22, cw - 0.55, 0.5)
    run(para(tf, first=True), title, size=16, color=WHITE, bold=True)
    tf = textbox(s, x + 0.3, y + 0.78, cw - 0.55, ch - 0.9)
    run(para(tf, first=True, line=1.05), body, size=12.5, color=MUTE)
page_no(s, 2)


# ── Slide 3 — Solution ────────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Solution")
tf = textbox(s, 0.7, 1.05, 11.9, 0.95)
p = para(tf, first=True)
run(p, "Rail Saarthi", size=34, color=WHITE, bold=True)
run(p, ": an agentic OS for Indian Railways", size=34, color=BLUE, bold=True)

cols = [
    ("What it is", "A multi-agent AI system: Train, Station, Crew, Maintenance and Passenger agents coordinating on a live digital twin.", BLUE),
    ("What it does", "Makes real-time decisions across trains, stations, crew and freight: detect a delay, reroute platforms, swap crew, alert passengers.", AMBER),
    ("What it replaces", "Reactive, manual control-room firefighting → proactive, autonomous coordination that self-heals the timetable.", GREEN),
]
cw, gx = 3.83, 0.2
x0, y = 0.7, 2.4
for i, (h, b, c) in enumerate(cols):
    x = x0 + i * (cw + gx)
    rect(s, x, y, cw, 2.95, fill=PANEL, line=LINE)
    rect(s, x, y, cw, 0.09, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, x + 0.32, y + 0.35, cw - 0.6, 0.6)
    run(para(tf, first=True), h, size=18, color=c, bold=True)
    tf = textbox(s, x + 0.32, y + 1.05, cw - 0.6, 1.7)
    run(para(tf, first=True, line=1.12), b, size=13.5, color=WHITE)

tf = textbox(s, 0.7, 5.75, 11.9, 0.7, align=PP_ALIGN.CENTER)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "Deterministic rules decide what's feasible. LLMs choose among safe options and explain why.",
    size=15, color=MUTE, italic=True)
page_no(s, 3)


# ── Slide 4 — Live demo ───────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Live Demo")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "Inject a 25-min delay. Watch the system self-heal.", size=30, color=WHITE, bold=True)

# screenshot frame (drop the dashboard capture here)
fx, fy, fw, fh = 0.7, 2.1, 8.2, 4.45
rect(s, fx, fy, fw, fh, fill=RGBColor(0x0D, 0x14, 0x26), line=LINE, line_w=1.5)
shot = BRAND / "dashboard.png"
if shot.exists():
    from PIL import Image as _Img

    iw, ih = _Img.open(shot).size
    aw, ah = fw - 0.16, fh - 0.16  # avail inside frame
    scale = min(aw / iw, ah / ih)
    dw, dh = iw * scale, ih * scale
    px = fx + 0.08 + (aw - dw) / 2
    py = fy + 0.08 + (ah - dh) / 2
    s.shapes.add_picture(str(shot), _in(px), _in(py), width=_in(dw), height=_in(dh))
else:
    tf = textbox(s, fx, fy + fh / 2 - 0.4, fw, 0.8, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    run(para(tf, first=True, align=PP_ALIGN.CENTER),
        "▶  Dashboard screenshot", size=18, color=MUTE, bold=True)
    run(para(tf, align=PP_ALIGN.CENTER, space_before=6),
        "railmind-demo.vercel.app  ·  drop a capture at frontend/public/brand/dashboard.png",
        size=11, color=LINE)

callouts = [
    ("Corridor map", "live train positions on the NDLS to DDU line"),
    ("Agent feed", "the cascade resolving, decision by decision"),
    ("Platform timeline", "a conflict caught and avoided"),
    ("KPI bar", "instant platforming, avg delay, knock-on saved"),
]
cx, cy = 9.1, 2.1
for t, b in callouts:
    rect(s, cx, cy, 3.5, 1.0, fill=PANEL, line=LINE)
    rect(s, cx, cy, 0.09, 1.0, fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, cx + 0.28, cy + 0.13, 3.1, 0.4)
    run(para(tf, first=True), t, size=14, color=BLUE, bold=True)
    tf = textbox(s, cx + 0.28, cy + 0.5, 3.1, 0.45)
    run(para(tf, first=True, line=1.0), b, size=11.5, color=MUTE)
    cy += 1.13
page_no(s, 4)


# ── Slide 5 — Passenger app ───────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Passenger App")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "The control room, in every passenger's pocket", size=30, color=WHITE, bold=True)

# phone mockup on the left
phone = BRAND / "passenger-app.png"
if phone.exists():
    from PIL import Image as _PImg

    iw, ih = _PImg.open(phone).size
    ph = 4.8  # target height in inches (clears the footer line)
    pw = ph * iw / ih
    s.shapes.add_picture(str(phone), _in(1.5), _in(2.05), width=_in(pw), height=_in(ph))

# feature callouts on the right
feats = [
    ("Live train status", "Delay, next stop and platform, refreshed in real time as the agents act.", AMBER),
    ("Proactive alerts", "Crew changes and disruptions are pushed to passengers before they even ask.", BLUE),
    ("Ask anything", "A voice or text assistant answers in plain language, in Hindi or English.", GREEN),
    ("Smart alternatives", "When a train runs late, it suggests the next service that gets you there.", AMBER),
]
cx, cy = 6.6, 2.15
for t, b, c in feats:
    rect(s, cx, cy, 6.0, 1.08, fill=PANEL, line=LINE)
    rect(s, cx, cy, 0.09, 1.08, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, cx + 0.32, cy + 0.16, 5.5, 0.4)
    run(para(tf, first=True), t, size=16, color=WHITE, bold=True)
    tf = textbox(s, cx + 0.32, cy + 0.56, 5.5, 0.45)
    run(para(tf, first=True, line=1.02), b, size=12.5, color=MUTE)
    cy += 1.2
page_no(s, 5)


# ── Slide 6 — Architecture ────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Architecture")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "Agents on a shared bus, fed by live rail data", size=30, color=WHITE, bold=True)

# agent row
agents = ["Train", "Station", "Crew", "Maintenance", "Passenger Info"]
aw, gx = 2.16, 0.18
x0, ay = 0.7, 2.25
for i, a in enumerate(agents):
    x = x0 + i * (aw + gx)
    rect(s, x, ay, aw, 0.85, fill=PANEL2, line=BLUE, line_w=1.25)
    tf = textbox(s, x, ay, aw, 0.85, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, a, size=14, color=WHITE, bold=True)
    run(para(tf, align=PP_ALIGN.CENTER), "Agent", size=9, color=MUTE, spacing=1)
    # connector down to bus
    rect(s, x + aw / 2 - 0.01, ay + 0.85, 0.025, 0.45, fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)

# message bus
by = 3.55
rect(s, 0.7, by, 11.9, 0.62, fill=BLUE, line=None)
tf = textbox(s, 0.7, by, 11.9, 0.62, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
run(para(tf, first=True, align=PP_ALIGN.CENTER), "CENTRAL MESSAGE BUS", size=15, color=NAVY, bold=True, spacing=3)

# data sources
tf = textbox(s, 0.7, 4.45, 11.9, 0.4)
run(para(tf, first=True), "DATA SOURCES", size=11, color=MUTE, bold=True, spacing=2)
sources = ["NTES", "IRCTC", "FOIS", "SCADA", "Weather"]
for i, src in enumerate(sources):
    x = x0 + i * (aw + gx)
    rect(s, x, 4.85, aw, 0.6, fill=PANEL, line=LINE)
    tf = textbox(s, x, 4.85, aw, 0.6, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), src, size=13, color=AMBER, bold=True)

rect(s, 0.7, 5.85, 11.9, 0.7, fill=PANEL, line=LINE)
tf = textbox(s, 0.7, 5.85, 11.9, 0.7, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "Stack:  ", size=13, color=MUTE, bold=True)
run(p, "FastAPI  ·  in-process event bus  ·  SQLite  ·  React  ·  Hive agents",
    size=13, color=WHITE)
page_no(s, 6)


# ── Slide 7 — Impact ──────────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Impact & Metrics")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "What the agents actually move", size=32, color=WHITE, bold=True)

metrics = [
    ("+34%", "network throughput", "vs naive dispatch (MDPI multi-agent study)", BLUE),
    ("90%", "instant platforming", "trains platformed without waiting, in sim", AMBER),
    ("<100ms", "agent response", "from event detected to decision proposed", GREEN),
    ("20%+", "fewer knock-on delays", "cascade reduction across the corridor", BLUE),
]
cw, ch, gx = 2.92, 2.9, 0.18
x0, y = 0.7, 2.25
for i, (big, label, sub, c) in enumerate(metrics):
    x = x0 + i * (cw + gx)
    rect(s, x, y, cw, ch, fill=PANEL, line=LINE)
    rect(s, x, y, cw, 0.09, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, x + 0.2, y + 0.5, cw - 0.4, 1.0, align=PP_ALIGN.CENTER)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), big, size=40, color=c, bold=True)
    tf = textbox(s, x + 0.2, y + 1.55, cw - 0.4, 0.5, align=PP_ALIGN.CENTER)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), label, size=15, color=WHITE, bold=True)
    tf = textbox(s, x + 0.2, y + 2.05, cw - 0.4, 0.7, align=PP_ALIGN.CENTER)
    run(para(tf, first=True, align=PP_ALIGN.CENTER, line=1.0), sub, size=11.5, color=MUTE)

tf = textbox(s, 0.7, 5.7, 11.9, 0.5, align=PP_ALIGN.CENTER)
run(para(tf, first=True, align=PP_ALIGN.CENTER),
    "Simulated results. Pilot validation planned.", size=13, color=MUTE, italic=True)
page_no(s, 7)


# ── Slide 8 — Roadmap ─────────────────────────────────────────────────────────
s = slide(prs)
chip(s, 0.7, 0.6, "Roadmap")
tf = textbox(s, 0.7, 1.05, 11.9, 0.9)
run(para(tf, first=True), "One corridor today, the network tomorrow", size=32, color=WHITE, bold=True)

phases = [
    ("Phase 1", "Hackathon MVP", ["Corridor simulation", "Delay injection", "Platform reassignment", "Live agent control room"], BLUE),
    ("Phase 2", "Pilot", ["5-station live corridor", "NTES integration", "Real crew & duty data", "Operator-in-the-loop"], AMBER),
    ("Phase 3", "Scale", ["National rollout", "RL-based dispatch", "Passenger AI chatbot", "Rail Madad integration"], GREEN),
]
cw, gx = 3.83, 0.2
x0, y = 0.7, 2.4
# timeline spine
rect(s, x0 + 0.2, y - 0.18, 11.5, 0.04, fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)
for i, (tag, name, items, c) in enumerate(phases):
    x = x0 + i * (cw + gx)
    # node
    rect(s, x + 0.2, y - 0.32, 0.28, 0.28, fill=c, line=None, shape=MSO_SHAPE.OVAL)
    rect(s, x, y + 0.1, cw, 3.5, fill=PANEL, line=LINE)
    rect(s, x, y + 0.1, cw, 0.09, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    tf = textbox(s, x + 0.32, y + 0.4, cw - 0.6, 0.4)
    run(para(tf, first=True), tag.upper(), size=12, color=c, bold=True, spacing=2)
    tf = textbox(s, x + 0.32, y + 0.78, cw - 0.6, 0.5)
    run(para(tf, first=True), name, size=20, color=WHITE, bold=True)
    tf = textbox(s, x + 0.32, y + 1.5, cw - 0.6, 2.0)
    for j, it in enumerate(items):
        p = para(tf, first=(j == 0), space_after=7)
        run(p, "•  ", size=13, color=c, bold=True)
        run(p, it, size=13.5, color=WHITE)
page_no(s, 8)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides)} slides)")
